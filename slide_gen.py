import copy
import math
from io import BytesIO
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Emu, Pt

from parser import BugData, ImageStep

TEMPLATE_PATH = Path(__file__).parent / "Template" / "New Layout.pptx"

# Image area position and size (inches) — from New Layout.pptx inspection
IMG_LEFT, IMG_TOP = 1.2467, 1.6255
IMG_W,    IMG_H   = 10.8399, 5.7217

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

_TYPE_COLORS = {
    "Current":   RGBColor(0xC0, 0x00, 0x00),
    "Reference": RGBColor(0xFF, 0xC0, 0x00),
    "Proposal":  RGBColor(0xFF, 0x66, 0x00),
}


def _in(value: float) -> Emu:
    return Inches(value)


def _find_layout(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0]


def _clone_slide_shapes(src_slide, dst_slide) -> None:
    """Replace dst_slide's shapes with deep copies of src_slide's shapes."""
    sp_tree = dst_slide.shapes._spTree
    # spTree layout: [0]=nvGrpSpPr, [1]=grpSpPr, [2+]=shapes from layout.
    # Remove all shapes add_slide() already added, keep only the two metadata nodes.
    for child in list(sp_tree)[2:]:
        sp_tree.remove(child)
    # Append shapes from template slide (skip its own metadata nodes).
    for el in list(src_slide.shapes._spTree)[2:]:
        sp_tree.append(copy.deepcopy(el))


def _remove_template_slides(prs: Presentation, count: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for _ in range(count):
        sld_id_lst.remove(sld_id_lst[0])


def _safe_set_text(tf, text: str) -> None:
    """Set the text of the first paragraph, preserving existing run formatting."""
    p = tf.paragraphs[0]
    if p.runs:
        for run in p.runs[1:]:
            run.text = ""
        p.runs[0].text = text
    else:
        p.text = text


def _set_type_label(tf, type_name: str, section_text: str, color: RGBColor) -> None:
    """Build type-label paragraph with two runs:
    - Run 1: "TypeName: " in type color, bold, 14pt
    - Run 2: description text in Calibri, black, 14pt, no bold/italic
    """
    _A = _A_NS
    p = tf.paragraphs[0]
    p_el = p._p

    for r_el in list(p_el.findall(f"{{{_A}}}r")):
        p_el.remove(r_el)

    def _append_run(text, rgb, bold, italic=False, font_name=None):
        r = p_el.makeelement(f"{{{_A}}}r", {})
        rPr = etree.SubElement(r, f"{{{_A}}}rPr")
        rPr.set("lang", "en-US")
        rPr.set("sz", "1400")  # 14pt in hundredths of a point
        rPr.set("b", "1" if bold else "0")
        rPr.set("i", "1" if italic else "0")
        sf = etree.SubElement(rPr, f"{{{_A}}}solidFill")
        sc = etree.SubElement(sf, f"{{{_A}}}srgbClr")
        sc.set("val", str(rgb))
        if font_name:
            lat = etree.SubElement(rPr, f"{{{_A}}}latin")
            lat.set("typeface", font_name)
        t_el = etree.SubElement(r, f"{{{_A}}}t")
        t_el.text = text
        end_rpr = p_el.find(f"{{{_A}}}endParaRPr")
        if end_rpr is not None:
            end_rpr.addprevious(r)
        else:
            p_el.append(r)

    label_part = f"{type_name}: " if section_text else f"{type_name}:"
    _append_run(label_part, color, bold=True)
    if section_text:
        _append_run(section_text, RGBColor(0, 0, 0), bold=False, italic=False, font_name="Calibri")


def _force_shape_autofit(shape) -> None:
    """Set spAutoFit directly in XML, replacing any conflicting autofit element."""
    bodyPr = shape.text_frame._txBody.find(f"{{{_A_NS}}}bodyPr")
    if bodyPr is None:
        return
    for tag in ("noAutofit", "normAutofit", "spAutoFit"):
        el = bodyPr.find(f"{{{_A_NS}}}{tag}")
        if el is not None:
            bodyPr.remove(el)
    bodyPr.append(bodyPr.makeelement(f"{{{_A_NS}}}spAutoFit", {}))
    shape.text_frame.word_wrap = True


def _fit_shape_height(shape, text: str, font_pt: float) -> None:
    """Explicitly set shape height to fit text, so the saved PPTX renders correctly
    without relying on PowerPoint's open-time recalculation."""
    width_pt = shape.width / 12700
    # Average proportional character width ≈ 0.52× font size
    chars_per_line = max(1, width_pt / (font_pt * 0.52))
    lines = sum(
        max(1, math.ceil(len(para) / chars_per_line)) if para else 1
        for para in (text or "").split("\n")
    )
    line_height_emu = int(font_pt * 1.4 * 12700)
    padding_emu = int(0.15 * 914400)
    shape.height = lines * line_height_emu + padding_emu


def _get_font_pt(shape, default: float = 12.0) -> float:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size / 12700
    return default


def _set_no_autofit(shape) -> None:
    """Remove all autofit elements so the shape keeps its explicitly set dimensions."""
    bodyPr = shape.text_frame._txBody.find(f"{{{_A_NS}}}bodyPr")
    if bodyPr is None:
        return
    for tag in ("noAutofit", "normAutofit", "spAutoFit"):
        el = bodyPr.find(f"{{{_A_NS}}}{tag}")
        if el is not None:
            bodyPr.remove(el)


def _fit_to_one_line(shape, text: str, default_pt: float, min_pt: float = 10.0) -> float:
    """Return the largest font size (≤ default_pt, ≥ min_pt) that fits text in one line."""
    if not text:
        return default_pt
    width_pt = shape.width / 12700
    max_pt = width_pt / (len(text) * 0.52)
    return max(min_pt, min(default_pt, max_pt))


def _set_text_font_size(tf, pt: float) -> None:
    """Set font size on every run in the text frame."""
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(pt)


def _set_no_line(shape) -> None:
    """Remove the visible border/outline from a shape by setting its line to noFill."""
    sp_pr = shape._element.spPr
    existing = sp_pr.find(f"{{{_A_NS}}}ln")
    if existing is not None:
        sp_pr.remove(existing)
    ln = etree.SubElement(sp_pr, f"{{{_A_NS}}}ln")
    etree.SubElement(ln, f"{{{_A_NS}}}noFill")


def _insert_image(slide, image_data: bytes, img_top: Emu | None = None) -> None:
    """Scale image to fit the main image area and add it to the slide.

    img_top overrides IMG_TOP when the header content has grown taller than
    the default layout, pushing the image down to avoid covering text.
    """
    if not image_data:
        return
    img = Image.open(BytesIO(image_data))
    img_w_px, img_h_px = img.size

    if img_top is None:
        img_top = _in(IMG_TOP)

    area_w = _in(IMG_W)
    # Shrink the available height by however much we pushed the image down.
    area_h = _in(IMG_H) - (img_top - _in(IMG_TOP))
    scale = min(area_w / img_w_px, area_h / img_h_px)
    final_w = int(img_w_px * scale)
    final_h = int(img_h_px * scale)

    left = _in(IMG_LEFT) + (area_w - final_w) // 2
    top  = img_top + (area_h - final_h) // 2

    slide.shapes.add_picture(BytesIO(image_data), left, top, final_w, final_h)


def _update_slide(
    slide,
    step: ImageStep,
    title: str,
    version_info: str,
    section_texts: dict[str, str],
) -> None:
    """Populate all text fields and insert the image on a cloned slide."""
    # Remove any existing picture shapes from the template (replaced by actual image)
    for shape in list(slide.shapes):
        if int(shape.shape_type) == 13:  # MSO_SHAPE_TYPE.PICTURE
            shape.element.getparent().remove(shape.element)

    title_shape = None
    type_label_shape = None

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        _set_no_line(shape)
        name = shape.name

        if name == "標題 2":
            title_shape = shape
            _safe_set_text(shape.text_frame, title)
            font_pt = _fit_to_one_line(shape, title, _get_font_pt(shape, 24.0))
            _set_text_font_size(shape.text_frame, font_pt)
            shape.text_frame.word_wrap = False
            _fit_shape_height(shape, title, font_pt)
            _set_no_autofit(shape)
            continue  # skip _force_shape_autofit — we own the dimensions

        elif name == "Title 1" and shape.top < _in(1.3):
            type_label_shape = shape
            section_text = section_texts.get(step.img_type, "")
            color = _TYPE_COLORS.get(step.img_type, RGBColor(0, 0, 0))
            _set_type_label(shape.text_frame, step.img_type, section_text, color)
            full_label = f"{step.img_type}: {section_text}" if section_text else f"{step.img_type}:"
            _fit_shape_height(shape, full_label, 14.0)

        elif name == "文字方塊 4":
            shape.top = Emu(0)
            _safe_set_text(shape.text_frame, "PX, fix in X/X")
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            for run in p.runs:
                run.font.size = Pt(16)

        elif name == "文字方塊 9":
            # Annotation box (intentionally off right edge of slide)
            _safe_set_text(shape.text_frame, step.text)
            _fit_shape_height(shape, step.text, _get_font_pt(shape, 12.0))

        _force_shape_autofit(shape)

    # Push the type label below the short description so they don't overlap.
    if title_shape and type_label_shape:
        gap = Emu(int(0.08 * 914400))  # 0.08 inch breathing room
        type_label_shape.top = title_shape.top + title_shape.height + gap

    # If the type label now sits below the default image top, move the image
    # down to match so text is never covered.
    img_top = _in(IMG_TOP)
    if type_label_shape:
        label_bottom = type_label_shape.top + type_label_shape.height
        gap = Emu(int(0.08 * 914400))
        if label_bottom + gap > img_top:
            img_top = label_bottom + gap

    _insert_image(slide, step.image_data, img_top)


def generate(bug_data: BugData, output_path: str) -> None:
    """Generate a PPTX at output_path — one slide per ImageStep."""
    prs = Presentation(str(TEMPLATE_PATH))
    n_template_slides = len(prs.slides)

    layout = _find_layout(prs, "1_只有標題 (no BK)")
    template_slide = prs.slides[0]

    steps = list(bug_data.image_steps)
    if not any(s.img_type == "Reference" for s in steps):
        insert_at = next(
            (i for i, s in enumerate(steps) if s.img_type == "Proposal"),
            len(steps),
        )
        steps.insert(insert_at, ImageStep(
            step_num=1, text="", img_type="Reference", image_file="", image_data=b"",
        ))

    for step in steps:
        new_slide = prs.slides.add_slide(layout)
        _clone_slide_shapes(template_slide, new_slide)
        _update_slide(new_slide, step, bug_data.title, bug_data.version_info, bug_data.section_texts)

    _remove_template_slides(prs, n_template_slides)
    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(steps)} slide(s))")
